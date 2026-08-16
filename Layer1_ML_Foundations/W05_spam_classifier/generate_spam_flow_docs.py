from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

base_dir = Path("c:/Users/viraj/Agentic AI Starts/viraj-agentic-ai-roadmap/Layer1_ML_Foundations/W05_spam_classifier")

doc = Document()
doc.add_heading("Email Spam Classifier", level=1)
doc.add_paragraph("Complete Flow Overview")
doc.add_paragraph("")
doc.add_paragraph("This project classifies SMS-like messages as spam or not spam using Python, TF-IDF, and Logistic Regression.")
doc.add_heading("1. Big Picture", level=2)
doc.add_paragraph("Input: raw SMS text")
doc.add_paragraph("Process: clean text, vectorize, train model, predict")
doc.add_paragraph("Output: Spam / Not Spam + confidence score")
doc.add_heading("2. End-to-End Flow", level=2)
doc.add_paragraph("data/spam.csv -> dataset with labeled messages")
doc.add_paragraph("src/data_preprocessing.py -> clean text and normalize labels")
doc.add_paragraph("src/train.py -> TF-IDF + Logistic Regression training")
doc.add_paragraph("models/spam_classifier.joblib -> saved trained model")
doc.add_paragraph("app/app.py -> Streamlit interface for user input")
doc.add_paragraph("src/predict.py -> classify new message and show confidence")
doc.add_heading("3. Training Pipeline", level=2)
doc.add_paragraph("1. Load dataset from CSV")
doc.add_paragraph("2. Clean text and remove noise")
doc.add_paragraph("3. Split data into train/test")
doc.add_paragraph("4. Convert text to TF-IDF vectors")
doc.add_paragraph("5. Train Logistic Regression model")
doc.add_paragraph("6. Evaluate accuracy and precision")
doc.add_paragraph("7. Save model to disk as .joblib")
doc.add_heading("4. Prediction Pipeline", level=2)
doc.add_paragraph("User enters message in Streamlit UI")
doc.add_paragraph("Text is cleaned using same preprocessing")
doc.add_paragraph("Saved TF-IDF vectorizer transforms text to numbers")
doc.add_paragraph("Model predicts Spam or Ham")
doc.add_paragraph("Confidence score is shown in the app")
doc.save(base_dir / "Spam_Classifier_Flow.docx")

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 245, 250)

box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Email Spam Classifier - Complete Flow"
p.alignment = PP_ALIGN.LEFT
p.runs[0].font.bold = True
p.runs[0].font.size = Pt(24)
p.runs[0].font.color.rgb = RGBColor(18, 45, 90)

box2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(1.2))
tf2 = box2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "From raw SMS data to spam / ham prediction with TF-IDF and Logistic Regression"
p2.alignment = PP_ALIGN.LEFT
p2.runs[0].font.size = Pt(20)
p2.runs[0].font.color.rgb = RGBColor(65, 65, 65)

box3 = slide.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.8))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(255, 255, 255)
box3.line.color.rgb = RGBColor(52, 152, 219)
box3.line.width = 2
box3_tf = box3.text_frame
box3_tf.word_wrap = True

lines = [
    "Data source: data/spam.csv",
    "Preprocessing: clean text + normalize labels",
    "Feature extraction: TF-IDF vectorizer",
    "Model training: Logistic Regression",
    "Evaluation: accuracy, precision, recall, F1",
    "Persistence: models/spam_classifier.joblib",
    "Prediction: user enters text in Streamlit app",
    "Output: Spam / Ham + confidence score",
]

for i, line in enumerate(lines):
    p3 = box3_tf.paragraphs[0] if i == 0 else box3_tf.add_paragraph()
    p3.text = line
    p3.runs[0].font.size = Pt(20)
    p3.runs[0].font.color.rgb = RGBColor(44, 62, 80)

prs.save(base_dir / "Spam_Classifier_Flow.pptx")

print("Created:", base_dir / "Spam_Classifier_Flow.docx")
print("Created:", base_dir / "Spam_Classifier_Flow.pptx")
